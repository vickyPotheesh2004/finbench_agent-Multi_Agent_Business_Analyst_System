"""
src/ingestion/pdf_ingestor.py

Production-Grade Multi-Format Financial Document Ingestor

Optimized for:
- FinanceBench
- SEC filings
- Windows
- Colab
- OCR fallback
- Large PDFs
- Multi-format ingestion
- Stable metadata extraction
- Table extraction
- Heading detection
"""

from __future__ import annotations

import csv
import json
import logging
import os
import re

from dataclasses import dataclass
from typing import Dict
from typing import List
from typing import Tuple

logger = logging.getLogger(__name__)

# ─────────────────────────────────────────────────────────────────────────────
# Constants
# ─────────────────────────────────────────────────────────────────────────────

SUPPORTED_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".doc",
    ".xlsx",
    ".xls",
    ".csv",
    ".pptx",
    ".html",
    ".htm",
    ".xml",
    ".json",
    ".txt",
}

MAX_TEXT_CHARS = 20_000_000

# FIX-v14 (2026-06-07): raised from 2000 -> 30000.
# FIX-v11 (text-strategy) + FIX-v12 (single-row tables) extract far more
# cells per page. The old 2000 cap filled up from early pages (cover, TOC,
# MD&A) and CHOPPED OFF the actual income statement / balance sheet, which
# in 10-Ks appears ~page 50-70. That deleted the revenue/PPE cells before
# extraction ran -> Q2/Q3 "missing input" failures. 30000 small dicts is
# ~6 MB, well within the 14 GB RAM budget (C4).
MAX_TABLES = 30000

# FIX (2026-06-21): structured_tables must NOT share the flat-cell cap.
# table_cells are INDIVIDUAL cells (tens of thousands); structured_tables are
# WHOLE tables (a few hundred per 10-K). The 30000-cell cap fills up from the
# front (cover/TOC/MD&A) and on big filings (CVS, JPMorgan, Nike) the actual
# balance sheet / income statement sits PAST the cutoff and was being
# truncated away entirely -> wrong-cell picks + RETRIEVAL_MISS. We keep ALL
# structured_tables (they drive the normalizer) with a generous separate cap.
MAX_STRUCTURED_TABLES = 20000

MAX_HEADINGS = 3000

HEADING_FONT_SIZE_MIN = 13.0

HTML_CHARS_PER_PAGE = 5000

# ─────────────────────────────────────────────────────────────────────────────
# Table Cell
# ─────────────────────────────────────────────────────────────────────────────


@dataclass(slots=True)
class TableCell:

    row_header: str = ""

    col_header: str = ""

    value: str = ""

    page: int = 0

    table_number: int = 0

    section: str = ""

    company: str = ""

    doc_type: str = ""

    fiscal_year: str = ""

    def to_dict(
        self,
    ) -> Dict:

        return {
            "row_header": self.row_header,
            "col_header": self.col_header,
            "value": self.value,
            "page": self.page,
            "table_number": self.table_number,
            "section": self.section,
            "company": self.company,
            "doc_type": self.doc_type,
            "fiscal_year": self.fiscal_year,
        }

# ─────────────────────────────────────────────────────────────────────────────
# PDF Ingestor
# ─────────────────────────────────────────────────────────────────────────────


class PDFIngestor:

    def __init__(
        self,
        enable_images: bool = False,
        llm_client=None,
    ):

        self.enable_images = enable_images

        self._llm = llm_client

    # ─────────────────────────────────────────────────────────────────────

    def run(
        self,
        state,
    ):

        doc_path = getattr(
            state,
            "document_path",
            "",
        ) or ""

        if not doc_path:

            logger.warning(
                "[N01] Missing document path"
            )

            return state

        if not os.path.exists(
            doc_path
        ):

            logger.error(
                "[N01] File not found: %s",
                doc_path,
            )

            return state

        result = self.ingest(
            doc_path
        )

        state.raw_text = result.get(
            "raw_text",
            "",
        )

        state.table_cells = result.get(
            "table_cells",
            [],
        )

        # Structure-preserving tables (2026-06-13) — additive; safe if the
        # state model lacks the field on older builds.
        try:
            state.structured_tables = result.get(
                "structured_tables",
                [],
            )
        except Exception:
            pass

        state.heading_positions = (
            result.get(
                "heading_positions",
                [],
            )
        )

        if not getattr(
            state,
            "company_name",
            "",
        ):

            state.company_name = (
                result.get(
                    "company_name",
                    "",
                )
            )

        if not getattr(
            state,
            "doc_type",
            "",
        ):

            state.doc_type = (
                result.get(
                    "doc_type",
                    "",
                )
            )

        if not getattr(
            state,
            "fiscal_year",
            "",
        ):

            state.fiscal_year = (
                result.get(
                    "fiscal_year",
                    "",
                )
            )

        logger.info(
            "[N01] chars=%d tables=%d headings=%d",
            len(state.raw_text),
            len(state.table_cells),
            len(state.heading_positions),
        )

        # Optional image processing

        if (
            self.enable_images
            and doc_path.lower().endswith(
                ".pdf"
            )
        ):

            try:

                from src.ingestion.image_processor import (
                    ImageProcessor,
                )

                processor = (
                    ImageProcessor(
                        enable_ocr=True,
                        enable_vision=self._llm
                        is not None,
                        llm_client=self._llm,
                    )
                )

                state = processor.run(
                    state
                )

            except Exception:

                logger.exception(
                    "[N01] Image processing failed"
                )

        return state

    # ─────────────────────────────────────────────────────────────────────

    def ingest(
        self,
        file_path: str,
    ) -> Dict:

        ext = os.path.splitext(
            file_path
        )[1].lower()

        if (
            ext
            not in SUPPORTED_EXTENSIONS
        ):

            logger.warning(
                "[N01] Unsupported extension: %s",
                ext,
            )

            return self._empty_result()

        try:

            if ext == ".pdf":

                return self._ingest_pdf(
                    file_path
                )

            if ext in (
                ".docx",
                ".doc",
            ):

                return self._ingest_docx(
                    file_path
                )

            if ext in (
                ".xlsx",
                ".xls",
            ):

                return self._ingest_xlsx(
                    file_path
                )

            if ext == ".csv":

                return self._ingest_csv(
                    file_path
                )

            if ext == ".pptx":

                return self._ingest_pptx(
                    file_path
                )

            if ext in (
                ".html",
                ".htm",
                ".xml",
            ):

                return self._ingest_html(
                    file_path
                )

            if ext == ".json":

                return self._ingest_json(
                    file_path
                )

            return self._ingest_txt(
                file_path
            )

        except Exception:

            logger.exception(
                "[N01] ingestion failed"
            )

            return self._empty_result()

    # ─────────────────────────────────────────────────────────────────────
    # PDF
    # ─────────────────────────────────────────────────────────────────────

    def _ingest_pdf(
        self,
        file_path: str,
    ) -> Dict:

        raw_text = ""

        table_cells = []

        structured_tables = []

        heading_positions = []

        # Text + Tables

        try:

            import pdfplumber

            with pdfplumber.open(
                file_path
            ) as pdf:

                for page_num, page in enumerate(
                    pdf.pages,
                    start=1,
                ):

                    try:

                        text = (
                            page.extract_text()
                            or ""
                        )

                    except Exception:

                        text = ""

                    raw_text += (
                        text + "\n"
                    )

                    # Tables
                    # FIX-v11 (2026-06-07): pdfplumber's default "lines"
                    # strategy needs ruled borders. Many 10-Ks (e.g.
                    # Activision) use BORDERLESS tables -> default returns
                    # nothing (cells=0). Retry with the "text" strategy,
                    # which infers columns from whitespace alignment.
                    # Zero new dependencies — pdfplumber already supports this.

                    try:

                        tables = (
                            page.extract_tables()
                            or []
                        )

                    except Exception:

                        tables = []

                    if not tables:

                        try:

                            tables = (
                                page.extract_tables(
                                    table_settings={
                                        "vertical_strategy": "text",
                                        "horizontal_strategy": "text",
                                        "snap_tolerance": 4,
                                        "join_tolerance": 4,
                                        "edge_min_length": 3,
                                        "min_words_vertical": 2,
                                        "min_words_horizontal": 1,
                                    }
                                )
                                or []
                            )

                            if tables:

                                logger.info(
                                    "[N01] page %d: text-strategy recovered %d tables",
                                    page_num,
                                    len(tables),
                                )

                        except Exception:

                            tables = []

                    for table_idx, table in enumerate(
                        tables,
                        start=1,
                    ):

                        # FIX-v12 (2026-06-07): only require >=1 row.
                        # The old `< 2` check dropped single-row tables —
                        # which on borderless 10-Ks (Activision) is MOST of
                        # them, causing cells=0 even though pdfplumber found
                        # 403 tables. Now we treat row 0 as data when there's
                        # only one row (no header row to peel off).
                        if not table:
                            continue

                        # Structure-preserving capture (2026-06-13):
                        # keep the FULL table shape before flattening, so
                        # downstream column/period-aware selection can tell
                        # FY2022 from FY2021 columns (fixes wrong-year picks).
                        # Additive — the flatten loop below is untouched.
                        try:
                            clean_rows = [
                                [
                                    (str(c).strip() if c else "")
                                    for c in row
                                ]
                                for row in table
                                if row
                            ]
                            if clean_rows:
                                structured_tables.append(
                                    {
                                        "page": page_num,
                                        "table_number": table_idx,
                                        "headers": clean_rows[0],
                                        "rows": (
                                            clean_rows[1:]
                                            if len(clean_rows) > 1
                                            else clean_rows
                                        ),
                                        "n_rows": len(clean_rows),
                                        "n_cols": max(
                                            (len(r) for r in clean_rows),
                                            default=0,
                                        ),
                                    }
                                )
                        except Exception:
                            pass

                        if len(table) >= 2:
                            headers = [
                                str(c).strip()
                                if c
                                else ""
                                for c in table[0]
                            ]
                            data_rows = table[1:]
                        else:
                            # Single-row table: no header, treat as data
                            headers = []
                            data_rows = table

                        for row in data_rows:

                            if not row:
                                continue

                            row_header = str(
                                row[0]
                            ).strip() if row[0] else ""

                            for col_idx, cell in enumerate(
                                row[1:],
                                start=1,
                            ):

                                value = (
                                    str(cell).strip()
                                    if cell
                                    else ""
                                )

                                if not value:
                                    continue

                                col_header = (
                                    headers[col_idx]
                                    if col_idx
                                    < len(headers)
                                    else ""
                                )

                                table_cells.append(
                                    TableCell(
                                        row_header=row_header,
                                        col_header=col_header,
                                        value=value,
                                        page=page_num,
                                        table_number=table_idx,
                                    ).to_dict()
                                )

        except Exception:

            logger.exception(
                "[N01] pdfplumber failed"
            )

        # Heading Extraction

        try:

            import fitz

            doc = fitz.open(
                file_path
            )

            for page_num, page in enumerate(
                doc,
                start=1,
            ):

                blocks = page.get_text(
                    "dict"
                ).get(
                    "blocks",
                    [],
                )

                for block in blocks:

                    for line in block.get(
                        "lines",
                        [],
                    ):

                        for span in line.get(
                            "spans",
                            [],
                        ):

                            text = (
                                span.get(
                                    "text",
                                    "",
                                ).strip()
                            )

                            if (
                                not text
                                or len(text)
                                < 3
                            ):

                                continue

                            font_size = float(
                                span.get(
                                    "size",
                                    0,
                                )
                            )

                            flags = int(
                                span.get(
                                    "flags",
                                    0,
                                )
                            )

                            is_bold = bool(
                                flags & 16
                            )

                            if (
                                font_size
                                >= HEADING_FONT_SIZE_MIN
                            ):

                                heading_positions.append(
                                    {
                                        "text": text,
                                        "font_size": round(
                                            font_size,
                                            1,
                                        ),
                                        "is_bold": is_bold,
                                        "page": page_num,
                                    }
                                )

            doc.close()

        except Exception:

            logger.exception(
                "[N01] heading extraction failed"
            )

        raw_text = raw_text[
            :MAX_TEXT_CHARS
        ]

        # OCR fallback

        if len(raw_text.strip()) < 1000:

            logger.warning(
                "[N01] Low text extraction — OCR fallback"
            )

            raw_text += (
                self._ocr_pdf(
                    file_path
                )
            )

        company, doc_type, fiscal_year = (
            self._extract_metadata(
                raw_text
            )
        )

        self._apply_metadata(
            table_cells,
            company,
            doc_type,
            fiscal_year,
        )

        return {
            "raw_text": raw_text,
            "table_cells": table_cells[
                :MAX_TABLES
            ],
            "structured_tables": structured_tables[
                :MAX_STRUCTURED_TABLES
            ],
            "heading_positions": self._dedupe_headings(
                heading_positions
            )[
                :MAX_HEADINGS
            ],
            "company_name": company,
            "doc_type": doc_type,
            "fiscal_year": fiscal_year,
        }

    # ─────────────────────────────────────────────────────────────────────
    # OCR
    # ─────────────────────────────────────────────────────────────────────

    def _ocr_pdf(
        self,
        file_path: str,
    ) -> str:

        text_output = ""

        try:

            import fitz
            import pytesseract

            from PIL import Image

            doc = fitz.open(
                file_path
            )

            for page in doc:

                pix = page.get_pixmap(
                    dpi=200
                )

                mode = (
                    "RGBA"
                    if pix.alpha
                    else "RGB"
                )

                image = Image.frombytes(
                    mode,
                    [
                        pix.width,
                        pix.height,
                    ],
                    pix.samples,
                )

                text = (
                    pytesseract.image_to_string(
                        image
                    )
                )

                text_output += (
                    text + "\n"
                )

            doc.close()

        except Exception:

            logger.exception(
                "[N01] OCR failed"
            )

        return text_output

    # ─────────────────────────────────────────────────────────────────────
    # DOCX
    # ─────────────────────────────────────────────────────────────────────

    def _ingest_docx(
        self,
        file_path: str,
    ) -> Dict:

        raw_text = ""

        heading_positions = []

        try:

            from docx import (
                Document,
            )

            doc = Document(
                file_path
            )

            for para in doc.paragraphs:

                text = para.text.strip()

                if text:

                    raw_text += (
                        text + "\n"
                    )

                style = (
                    para.style.name
                    if para.style
                    else ""
                )

                if (
                    "Heading"
                    in style
                ):

                    heading_positions.append(
                        {
                            "text": text,
                            "font_size": 16.0,
                            "is_bold": True,
                            "page": 0,
                        }
                    )

        except Exception:

            logger.exception(
                "[N01] DOCX failed"
            )

        return {
            "raw_text": raw_text,
            "table_cells": [],
            "heading_positions": heading_positions,
            "company_name": "",
            "doc_type": "",
            "fiscal_year": "",
        }

    # ─────────────────────────────────────────────────────────────────────
    # XLSX
    # ─────────────────────────────────────────────────────────────────────

    def _ingest_xlsx(
        self,
        file_path: str,
    ) -> Dict:

        raw_text = ""

        try:

            import openpyxl

            wb = (
                openpyxl.load_workbook(
                    file_path,
                    data_only=True,
                )
            )

            for sheet in wb.sheetnames:

                ws = wb[sheet]

                for row in ws.iter_rows(
                    values_only=True
                ):

                    values = [
                        str(v).strip()
                        for v in row
                        if v is not None
                    ]

                    if values:

                        raw_text += (
                            " | ".join(values)
                            + "\n"
                        )

        except Exception:

            logger.exception(
                "[N01] XLSX failed"
            )

        return {
            "raw_text": raw_text,
            "table_cells": [],
            "heading_positions": [],
            "company_name": "",
            "doc_type": "",
            "fiscal_year": "",
        }

    # ─────────────────────────────────────────────────────────────────────
    # CSV
    # ─────────────────────────────────────────────────────────────────────

    def _ingest_csv(
        self,
        file_path: str,
    ) -> Dict:

        raw_text = ""

        try:

            with open(
                file_path,
                "r",
                encoding="utf-8",
                errors="replace",
            ) as f:

                reader = csv.reader(
                    f
                )

                for row in reader:

                    values = [
                        str(v).strip()
                        for v in row
                    ]

                    raw_text += (
                        " | ".join(values)
                        + "\n"
                    )

        except Exception:

            logger.exception(
                "[N01] CSV failed"
            )

        return {
            "raw_text": raw_text,
            "table_cells": [],
            "heading_positions": [],
            "company_name": "",
            "doc_type": "",
            "fiscal_year": "",
        }

    # ─────────────────────────────────────────────────────────────────────
    # PPTX
    # ─────────────────────────────────────────────────────────────────────

    def _ingest_pptx(
        self,
        file_path: str,
    ) -> Dict:

        raw_text = ""

        try:

            from pptx import (
                Presentation,
            )

            prs = Presentation(
                file_path
            )

            for slide in prs.slides:

                for shape in slide.shapes:

                    if (
                        hasattr(
                            shape,
                            "text",
                        )
                    ):

                        text = (
                            shape.text.strip()
                        )

                        if text:

                            raw_text += (
                                text + "\n"
                            )

        except Exception:

            logger.exception(
                "[N01] PPTX failed"
            )

        return {
            "raw_text": raw_text,
            "table_cells": [],
            "heading_positions": [],
            "company_name": "",
            "doc_type": "",
            "fiscal_year": "",
        }

    # ─────────────────────────────────────────────────────────────────────
    # HTML
    # ─────────────────────────────────────────────────────────────────────

    def _ingest_html(
        self,
        file_path: str,
    ) -> Dict:

        try:

            from bs4 import (
                BeautifulSoup,
            )

            with open(
                file_path,
                "r",
                encoding="utf-8",
                errors="replace",
            ) as f:

                html = f.read()

            soup = BeautifulSoup(
                html,
                "lxml",
            )

            raw_text = soup.get_text(
                separator="\n"
            )

            raw_text = re.sub(
                r"\n{3,}",
                "\n\n",
                raw_text,
            )

        except Exception:

            logger.exception(
                "[N01] HTML failed"
            )

            raw_text = ""

        return {
            "raw_text": raw_text,
            "table_cells": [],
            "heading_positions": [],
            "company_name": "",
            "doc_type": "",
            "fiscal_year": "",
        }

    # ─────────────────────────────────────────────────────────────────────
    # TXT
    # ─────────────────────────────────────────────────────────────────────

    def _ingest_txt(
        self,
        file_path: str,
    ) -> Dict:

        try:

            with open(
                file_path,
                "r",
                encoding="utf-8",
                errors="replace",
            ) as f:

                raw_text = f.read()

        except Exception:

            logger.exception(
                "[N01] TXT failed"
            )

            raw_text = ""

        return {
            "raw_text": raw_text,
            "table_cells": [],
            "heading_positions": [],
            "company_name": "",
            "doc_type": "",
            "fiscal_year": "",
        }

    # ─────────────────────────────────────────────────────────────────────
    # JSON
    # ─────────────────────────────────────────────────────────────────────

    def _ingest_json(
        self,
        file_path: str,
    ) -> Dict:

        try:

            with open(
                file_path,
                "r",
                encoding="utf-8",
                errors="replace",
            ) as f:

                data = json.load(
                    f
                )

            raw_text = json.dumps(
                data,
                indent=2,
            )

        except Exception:

            logger.exception(
                "[N01] JSON failed"
            )

            raw_text = ""

        return {
            "raw_text": raw_text,
            "table_cells": [],
            "heading_positions": [],
            "company_name": "",
            "doc_type": "",
            "fiscal_year": "",
        }

    # ─────────────────────────────────────────────────────────────────────
    # Metadata
    # ─────────────────────────────────────────────────────────────────────

    def _extract_metadata(
        self,
        text: str,
    ) -> Tuple[
        str,
        str,
        str,
    ]:

        return (
            self._extract_company(
                text
            ),
            self._extract_doc_type(
                text
            ),
            self._extract_fiscal_year(
                text
            ),
        )

    @staticmethod
    def _extract_company(
        text: str,
    ) -> str:

        snippet = text[:5000]

        # FIX (2026-06-21): the old single greedy regex matched the FIRST
        # "...Inc|Corp|LLC" phrase in the cover page, which is often exchange
        # boilerplate from the "securities registered" block (e.g. "The NASDAQ
        # Stock Market LLC", "New York Stock Exchange LLC") rather than the
        # filer. Two improvements:
        #   1. Prefer the registrant name that 10-K cover pages print right
        #      before "(Exact name of registrant as specified in its charter)".
        #   2. Reject known exchange / boilerplate phrases.
        _BOILERPLATE = (
            "nasdaq", "new york stock exchange", "nyse",
            "stock market", "stock exchange", "securities exchange",
            "the depository trust", "cboe",
        )

        def _is_boilerplate(name: str) -> bool:
            low = name.lower()
            return any(b in low for b in _BOILERPLATE)

        # 1. Registrant line: "<NAME>\n(Exact name of registrant ...)"
        reg = re.search(
            r"([A-Z][A-Za-z0-9&.,\-\s]{2,80}?)\s*\n?\s*\(\s*Exact name of",
            snippet,
            re.I,
        )
        if reg:
            cand = reg.group(1).strip().strip(",").strip()
            # keep only the last line (the name sits just above the marker)
            cand = cand.splitlines()[-1].strip() if cand else cand
            if cand and not _is_boilerplate(cand):
                return cand

        # 2. Fallback: first corporate-suffix phrase that is NOT boilerplate.
        pattern = (
            r"((?:[A-Z][a-zA-Z&,\-]+\s*){1,6}"
            r"(?:Inc|Corp|Corporation|Ltd|LLC|Company|Holdings|Group)\.?)"
        )
        for match in re.finditer(pattern, snippet):
            cand = match.group(1).strip()
            if cand and not _is_boilerplate(cand):
                return cand

        return ""

    @staticmethod
    def _extract_doc_type(
        text: str,
    ) -> str:

        upper = text.upper()

        for dtype in (
            "10-K",
            "10-Q",
            "8-K",
            "20-F",
            "6-K",
            "DEF 14A",
        ):

            if (
                dtype.replace(
                    "-",
                    "",
                )
                in upper.replace(
                    "-",
                    "",
                )
            ):

                return dtype

        return "UNKNOWN"

    @staticmethod
    def _extract_fiscal_year(
        text: str,
    ) -> str:

        patterns = [
            r"FY\s*(\d{4})",
            r"Fiscal Year.*?(\d{4})",
            r"Year Ended.*?(\d{4})",
        ]

        for pattern in patterns:

            match = re.search(
                pattern,
                text[:10000],
                re.I | re.S,
            )

            if match:

                return (
                    f"FY{match.group(1)}"
                )

        return ""

    # ─────────────────────────────────────────────────────────────────────
    # Helpers
    # ─────────────────────────────────────────────────────────────────────

    @staticmethod
    def _dedupe_headings(
        headings: List[Dict],
    ) -> List[Dict]:

        seen = set()

        output = []

        for h in headings:

            key = (
                h.get(
                    "text",
                    "",
                )
                .strip()
                .lower(),
                int(
                    h.get(
                        "page",
                        0,
                    )
                ),
            )

            if (
                not key[0]
                or key in seen
            ):

                continue

            seen.add(key)

            output.append(h)

        return output

    @staticmethod
    def _apply_metadata(
        table_cells: List[Dict],
        company_name: str,
        doc_type: str,
        fiscal_year: str,
    ):

        for cell in table_cells:

            cell["company"] = (
                company_name
                or "UNKNOWN"
            )

            cell["doc_type"] = (
                doc_type
                or "UNKNOWN"
            )

            cell["fiscal_year"] = (
                fiscal_year
                or "UNKNOWN"
            )

    @staticmethod
    def _empty_result():

        return {
            "raw_text": "",
            "table_cells": [],
            "heading_positions": [],
            "company_name": "",
            "doc_type": "",
            "fiscal_year": "",
        }

# ─────────────────────────────────────────────────────────────────────────────
# Convenience Wrapper
# ─────────────────────────────────────────────────────────────────────────────


def run_pdf_ingestor(
    state,
    enable_images: bool = False,
    llm_client=None,
):

    ingestor = PDFIngestor(
        enable_images=enable_images,
        llm_client=llm_client,
    )

    return ingestor.run(
        state
    )